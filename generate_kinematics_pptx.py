#!/usr/bin/env python3
"""Generate '1D Kinematics Intro' PowerPoint deck (16:9, modern high-contrast design)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)       # dark navy - headers / title bg
NAVY_LIGHT = RGBColor(0x14, 0x2C, 0x52)  # slightly lighter navy for gradients/accents
ACCENT = RGBColor(0x2E, 0xC4, 0xB6)      # teal accent
ACCENT_DARK = RGBColor(0x1E, 0x8F, 0x84)
GOLD = RGBColor(0xF2, 0xA6, 0x3D)        # warm accent for highlights
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF7, 0xF9, 0xFC)
LIGHT_GRAY = RGBColor(0xE9, 0xED, 0xF3)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA3)
TEXT_DARK = RGBColor(0x1B, 0x22, 0x2D)
TEXT_BODY = RGBColor(0x33, 0x3D, 0x4C)
CARD_BLUE = RGBColor(0xEA, 0xF3, 0xF7)
CARD_BORDER = RGBColor(0x2E, 0xC4, 0xB6)
RED_ACCENT = RGBColor(0xE0, 0x5A, 0x47)

FONT = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False, line_color=None, line_w=Pt(1), shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or color
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        effectLst = el.makeelement(qn('a:effectLst'), {})
        outerShdw = el.makeelement(qn('a:outerShdw'), {
            'blurRad': '90000', 'dist': '30000', 'dir': '5400000', 'rotWithShape': '0'
        })
        clr = el.makeelement(qn('a:srgbClr'), {'val': '1B222D'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '28000'})
        clr.append(alpha)
        outerShdw.append(clr)
        effectLst.append(outerShdw)
        el.append(effectLst)
    return shp


def add_rounded_rect(slide, x, y, w, h, color, line=False, line_color=None, line_w=Pt(1.25), radius=0.06, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or color
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        effectLst = el.makeelement(qn('a:effectLst'), {})
        outerShdw = el.makeelement(qn('a:outerShdw'), {
            'blurRad': '90000', 'dist': '30000', 'dir': '5400000', 'rotWithShape': '0'
        })
        clr = el.makeelement(qn('a:srgbClr'), {'val': '1B222D'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(alpha)
        outerShdw.append(clr)
        effectLst.append(outerShdw)
        el.append(effectLst)
    return shp


def add_textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return run


def add_para(tf, first=False, align=PP_ALIGN.LEFT, space_after=6, space_before=0, line_spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = line_spacing
    return p


def add_bullet_marker(slide, x, y, size, color):
    """Small square bullet marker."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header_bar(slide, kicker, title, slide_num, total=10):
    """Standard navy header bar used on content slides."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), NAVY)
    add_rect(slide, 0, Inches(1.35), SLIDE_W, Pt(3.5), ACCENT)
    # kicker
    tb, tf = add_textbox(slide, Inches(0.6), Inches(0.18), Inches(8), Inches(0.35))
    p = add_para(tf, first=True)
    set_run(p.add_run(), kicker.upper(), 12, ACCENT, bold=True)
    p.runs[0].font._rPr.set(qn('spc'), '150') if False else None
    # title
    tb2, tf2 = add_textbox(slide, Inches(0.6), Inches(0.5), Inches(10.5), Inches(0.8))
    p2 = add_para(tf2, first=True)
    set_run(p2.add_run(), title, 30, WHITE, bold=True)
    # slide number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.55), Inches(0.42), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge.shadow.inherit = False
    btf = badge.text_frame
    btf.margin_left = 0; btf.margin_right = 0; btf.margin_top = 0; btf.margin_bottom = 0
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    set_run(bp.add_run(), str(slide_num), 16, NAVY, bold=True)


def footer(slide, slide_num):
    tb, tf = add_textbox(slide, Inches(0.6), Inches(7.14), Inches(6), Inches(0.3))
    p = add_para(tf, first=True)
    set_run(p.add_run(), "Physics Mechanics 101  |  1D Kinematics", 10, MID_GRAY)
    tb2, tf2 = add_textbox(slide, Inches(11.8), Inches(7.14), Inches(1.2), Inches(0.3))
    p2 = add_para(tf2, first=True, align=PP_ALIGN.RIGHT)
    set_run(p2.add_run(), f"{slide_num} / 10", 10, MID_GRAY)


def bullet_list(slide, x, y, w, items, size=16, color=TEXT_BODY, gap=14, marker_color=ACCENT,
                 bold_terms=None, line_spacing=1.08):
    """items: list of (text) or (text, sub_items list) for two-level bullets.
    bold_terms: list of substrings to bold within each line (simple, non-overlapping)."""
    cur_y = y
    for item in items:
        if isinstance(item, tuple):
            if len(item) == 1:
                text, subitems = item[0], None
            else:
                text, subitems = item
        else:
            text, subitems = item, None
        line_h = Inches(0.05 + size / 72 * line_spacing)
        marker_size = Emu(int(Pt(size * 0.34)))
        marker_y = cur_y + Emu(int(Pt(size * 0.28)))
        add_bullet_marker(slide, x, marker_y, marker_size, marker_color)
        tb, tf = add_textbox(slide, x + Inches(0.28), cur_y, w - Inches(0.28), Inches(1.2))
        p = add_para(tf, first=True, space_after=2, line_spacing=line_spacing)
        _add_rich_line(p, text, size, color, bold_terms)
        cur_y += Inches((size / 72 * line_spacing) + 0.14)
        if subitems:
            for sub in subitems:
                stb, stf = add_textbox(slide, x + Inches(0.62), cur_y, w - Inches(0.62), Inches(1.0))
                sp = add_para(stf, first=True, space_after=2, line_spacing=line_spacing)
                set_run(sp.add_run(), "– ", size - 2, MID_GRAY)
                _add_rich_line(sp, sub, size - 2, TEXT_BODY, bold_terms, append=True)
                cur_y += Inches(((size - 2) / 72 * line_spacing) + 0.1)
        cur_y += Inches(gap / 72)
    return cur_y


def _add_rich_line(p, text, size, color, bold_terms, append=False):
    """Add text to paragraph p, bolding any bold_terms substrings found."""
    bold_terms = bold_terms or []
    remaining = text
    # Simple approach: split on ** markers if author used them, else bold listed terms
    if "**" in text:
        parts = text.split("**")
        for i, part in enumerate(parts):
            if not part:
                continue
            r = p.add_run()
            set_run(r, part, size, color, bold=(i % 2 == 1))
    else:
        r = p.add_run()
        set_run(r, text, size, color, bold=False)


def formula_card(slide, x, y, w, h, label, formula, note=None, accent=ACCENT, bg=CARD_BLUE):
    card = add_rounded_rect(slide, x, y, w, h, bg, line=True, line_color=accent, line_w=Pt(1.5), radius=0.12)
    accent_bar = add_rect(slide, x, y, Inches(0.09), h, accent)
    pad = Inches(0.28)
    if label:
        tb, tf = add_textbox(slide, x + pad, y + Inches(0.15), w - pad * 2, Inches(0.3))
        p = add_para(tf, first=True)
        set_run(p.add_run(), label.upper(), 11.5, ACCENT_DARK, bold=True)
        formula_y = y + Inches(0.48)
    else:
        formula_y = y + Inches(0.18)
    remaining_h = Inches(max(0.4, h.inches - Emu(formula_y - y).inches - 0.15))
    tb2, tf2 = add_textbox(slide, x + pad, formula_y, w - pad * 2, remaining_h)
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE if not note else MSO_ANCHOR.TOP
    p2 = add_para(tf2, first=True, align=PP_ALIGN.LEFT)
    set_run(p2.add_run(), formula, 24, NAVY, bold=True, font="Cambria")
    if note:
        p3 = add_para(tf2, space_before=6)
        set_run(p3.add_run(), note, 13, TEXT_BODY)
    return card


def section_label(slide, x, y, text, color=NAVY):
    tb, tf = add_textbox(slide, x, y, Inches(6), Inches(0.4))
    p = add_para(tf, first=True)
    set_run(p.add_run(), text, 16, color, bold=True)
    return tb


def add_table(slide, x, y, w, h, rows_data, col_widths, header_bg=NAVY, header_color=WHITE,
              body_size=14, header_size=14, row_colors=(WHITE, OFFWHITE)):
    rows = len(rows_data)
    cols = len(rows_data[0])
    gshape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = gshape.table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.15)
            cell.margin_right = Inches(0.15)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.08)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            if r == 0:
                set_run(run, val, header_size, header_color, bold=True)
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                bold_first = (c == 0)
                set_run(run, val, body_size, TEXT_DARK, bold=bold_first)
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_colors[r % 2]
    # remove default table style banding by keeping our own fills (already applied)
    return table


def pill(slide, x, y, w, h, text, bg, color, size=13):
    shp = add_rounded_rect(slide, x, y, w, h, bg, radius=0.5)
    tf = shp.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, size, color, bold=True)
    return shp


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, NAVY)
# decorative diagonal accent band
band = slide_shapes = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(-1.5), Inches(5.3), Inches(17), Inches(3.2))
band.rotation = 0
band.fill.solid()
band.fill.fore_color.rgb = NAVY_LIGHT
band.line.fill.background()
band.shadow.inherit = False

add_rect(s, 0, Inches(6.85), SLIDE_W, Inches(0.06), ACCENT)

# small motion-arrow motif
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.9), Inches(1.05), Inches(1.3), Inches(0.28))
arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT; arrow.line.fill.background(); arrow.shadow.inherit = False

tb, tf = add_textbox(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.5))
p = add_para(tf, first=True)
set_run(p.add_run(), "PHYSICS MECHANICS 101", 20, ACCENT, bold=True)

tb2, tf2 = add_textbox(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(2.0))
p2 = add_para(tf2, first=True, line_spacing=1.05)
set_run(p2.add_run(), "Introduction to", 40, WHITE, bold=False)
p3 = add_para(tf2, line_spacing=1.05)
set_run(p3.add_run(), "1D Kinematics", 54, WHITE, bold=True)

tb3, tf3 = add_textbox(s, Inches(0.9), Inches(4.15), Inches(9.5), Inches(0.6))
p4 = add_para(tf3, first=True)
set_run(p4.add_run(), "Describing motion along a straight line — position, velocity, acceleration, and the equations that connect them.", 16, LIGHT_GRAY)

# bottom info strip
tb4, tf4 = add_textbox(s, Inches(0.9), Inches(6.95), Inches(8), Inches(0.4))
p5 = add_para(tf4, first=True)
set_run(p5.add_run(), "Unit 1  •  Kinematics in One Dimension", 13, RGBColor(0xCF, 0xD8, 0xE6), bold=True)


# ---------------------------------------------------------------------------
# SLIDE 2 — What is Kinematics?
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Foundations", "What is Kinematics?", 2)

# Definition callout
def_card = add_rounded_rect(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(1.15), NAVY, radius=0.1)
tb, tf = add_textbox(s, Inches(0.95), Inches(1.8), Inches(11.4), Inches(0.9), anchor=MSO_ANCHOR.MIDDLE)
p = add_para(tf, first=True, line_spacing=1.15)
set_run(p.add_run(), "Kinematics", 19, ACCENT, bold=True)
set_run(p.add_run(), "  is the branch of mechanics that describes ", 17, WHITE)
set_run(p.add_run(), "how objects move", 17, GOLD, bold=True)
set_run(p.add_run(), " — without asking what causes the motion (that's ", 17, WHITE)
set_run(p.add_run(), "dynamics", 17, ACCENT, bold=True)
set_run(p.add_run(), ").", 17, WHITE)

# Two comparison cards: Scalar vs Vector
card_y = Inches(3.05)
card_h = Inches(2.15)
card_w = Inches(5.85)

c1 = add_rounded_rect(s, Inches(0.6), card_y, card_w, card_h, OFFWHITE, line=True, line_color=LIGHT_GRAY, line_w=Pt(1), radius=0.07, shadow=True)
tb, tf = add_textbox(s, Inches(0.9), card_y + Inches(0.2), card_w - Inches(0.6), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "🔵 SCALAR QUANTITIES", 15, ACCENT_DARK, bold=True)
bullet_list(s, Inches(0.9), card_y + Inches(0.65), card_w - Inches(0.6),
            [("Magnitude **only** — just a number + unit",),
             ("Examples: **distance**, **speed**, time, mass",),
             ("e.g. “5 meters” or “10 m/s”",)],
            size=14.5, gap=8)

c2 = add_rounded_rect(s, Inches(6.85), card_y, card_w, card_h, OFFWHITE, line=True, line_color=LIGHT_GRAY, line_w=Pt(1), radius=0.07, shadow=True)
tb, tf = add_textbox(s, Inches(7.15), card_y + Inches(0.2), card_w - Inches(0.6), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "🟠 VECTOR QUANTITIES", 15, RGBColor(0xC7, 0x6A, 0x1B), bold=True)
bullet_list(s, Inches(7.15), card_y + Inches(0.65), card_w - Inches(0.6),
            [("Magnitude **and** direction",),
             ("Examples: **displacement**, **velocity**, acceleration",),
             ("e.g. “5 meters, East” or “10 m/s, +x”",)],
            size=14.5, gap=8)

# 1D motion note
note = add_rounded_rect(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.15), CARD_BLUE, line=True, line_color=ACCENT, line_w=Pt(1.25), radius=0.1)
tb, tf = add_textbox(s, Inches(0.95), Inches(5.6), Inches(11.4), Inches(0.9), anchor=MSO_ANCHOR.MIDDLE)
p = add_para(tf, first=True, line_spacing=1.1)
set_run(p.add_run(), "Motion in 1D:  ", 15.5, NAVY, bold=True)
set_run(p.add_run(), "In this unit, all motion happens along a ", 15, TEXT_BODY)
set_run(p.add_run(), "single straight-line axis", 15, TEXT_BODY, bold=True)
set_run(p.add_run(), " (usually the x-axis). Direction is shown simply with a ", 15, TEXT_BODY)
set_run(p.add_run(), "positive (+) or negative (–) sign", 15, TEXT_BODY, bold=True)
set_run(p.add_run(), ".", 15, TEXT_BODY)

footer(s, 2)


# ---------------------------------------------------------------------------
# SLIDE 3 — Position & Displacement
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Core Concepts", "Position & Displacement", 3)

section_label(s, Inches(0.6), Inches(1.65), "Position (x)")
bullet_list(s, Inches(0.6), Inches(2.1), Inches(5.8),
            [("An object's **location** relative to a chosen origin (reference point)",),
             ("Measured in **meters (m)**",),
             ("Can be positive or negative depending on which side of the origin",),
             ],
            size=15.5, gap=12)

section_label(s, Inches(0.6), Inches(4.35), "Displacement (Δx)")
bullet_list(s, Inches(0.6), Inches(4.8), Inches(5.8),
            [("The **change in position** — how far out of place an object is",),
             ("A **vector**: only depends on start & end points, not the path taken",),
             ("Can be zero even if distance traveled is large (round trip)",),
             ],
            size=15.5, gap=12)

# formula card on right
formula_card(s, Inches(6.85), Inches(1.75), Inches(5.85), Inches(2.0),
             "Displacement Formula",
             "Δx = x_f − x_i",
             note="x_f = final position     x_i = initial position", accent=ACCENT, bg=CARD_BLUE)

# mini number-line visual
line_y = Inches(4.35)
add_rect(s, Inches(6.95), line_y + Inches(0.85), Inches(5.6), Pt(2.5), MID_GRAY)
# arrowhead
tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(12.45), line_y + Inches(0.72), Inches(0.22), Inches(0.28))
tri.rotation = 90
tri.fill.solid(); tri.fill.fore_color.rgb = MID_GRAY; tri.line.fill.background(); tri.shadow.inherit = False

xi_dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.6), line_y + Inches(0.72), Inches(0.28), Inches(0.28))
xi_dot.fill.solid(); xi_dot.fill.fore_color.rgb = GOLD; xi_dot.line.color.rgb = WHITE; xi_dot.line.width = Pt(1.5); xi_dot.shadow.inherit = False
xf_dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.1), line_y + Inches(0.72), Inches(0.28), Inches(0.28))
xf_dot.fill.solid(); xf_dot.fill.fore_color.rgb = ACCENT_DARK; xf_dot.line.color.rgb = WHITE; xf_dot.line.width = Pt(1.5); xf_dot.shadow.inherit = False

tb, tf = add_textbox(s, Inches(7.35), line_y + Inches(1.05), Inches(1.0), Inches(0.3))
p = add_para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "xᵢ", 13, GOLD, bold=True)
tb, tf = add_textbox(s, Inches(10.85), line_y + Inches(1.05), Inches(1.0), Inches(0.3))
p = add_para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "x_f", 13, ACCENT_DARK, bold=True)

tb, tf = add_textbox(s, Inches(6.95), line_y + Inches(0.2), Inches(5.6), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "Origin → direction of travel along the x-axis", 13, MID_GRAY, italic=True)

# arrow showing displacement above dots
disp_arrow = s.shapes.add_connector(2, Inches(7.74), line_y + Inches(0.6), Inches(11.24), line_y + Inches(0.6))
disp_arrow.line.color.rgb = NAVY
disp_arrow.line.width = Pt(1.5)

footer(s, 3)


# ---------------------------------------------------------------------------
# SLIDE 4 — Distance vs. Displacement
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Key Comparison", "Distance vs. Displacement", 4)

tb, tf = add_textbox(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "Two quantities that sound similar but behave very differently.", 15, TEXT_BODY, italic=True)

rows = [
    ["", "Distance", "Displacement"],
    ["Quantity type", "Scalar", "Vector"],
    ["Definition", "Total length of the path traveled", "Straight-line change from start to end"],
    ["Direction?", "No", "Yes (+ / –)"],
    ["Symbol", "d", "Δx"],
    ["Units (SI)", "meters (m)", "meters (m)"],
    ["Can it be zero\nafter moving?", "No (always ≥ 0)", "Yes, if start = end point"],
]
col_widths = [Inches(2.6), Inches(4.75), Inches(4.75)]
tbl = add_table(s, Inches(0.6), Inches(2.05), Inches(12.1), Inches(3.6), rows, col_widths,
                 header_bg=NAVY, header_color=WHITE, body_size=14.5, header_size=15.5)
# style first column of body rows
for r in range(1, len(rows)):
    cell = tbl.cell(r, 0)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY_LIGHT if False else LIGHT_GRAY
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = NAVY
            run.font.bold = True

# Example callout
ex = add_rounded_rect(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.95), CARD_BLUE, line=True, line_color=ACCENT, radius=0.12)
tb, tf = add_textbox(s, Inches(0.95), Inches(5.98), Inches(11.4), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
p = add_para(tf, first=True)
set_run(p.add_run(), "Example:  ", 14.5, NAVY, bold=True)
set_run(p.add_run(), "Walk 8 m east, then 3 m west → ", 14.5, TEXT_BODY)
set_run(p.add_run(), "Distance = 11 m", 14.5, RED_ACCENT, bold=True)
set_run(p.add_run(), "   but   ", 14.5, TEXT_BODY)
set_run(p.add_run(), "Displacement = 5 m East", 14.5, ACCENT_DARK, bold=True)

footer(s, 4)


# ---------------------------------------------------------------------------
# SLIDE 5 — Speed & Velocity
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Core Concepts", "Speed & Velocity", 5)

card_y = Inches(1.7)
card_h = Inches(2.35)
card_w = Inches(5.85)

c1 = add_rounded_rect(s, Inches(0.6), card_y, card_w, card_h, OFFWHITE, line=True, line_color=LIGHT_GRAY, radius=0.08, shadow=True)
tb, tf = add_textbox(s, Inches(0.9), card_y + Inches(0.2), card_w - Inches(0.6), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "AVERAGE", 14, ACCENT_DARK, bold=True)
bullet_list(s, Inches(0.9), card_y + Inches(0.62), card_w - Inches(0.6),
            [("**Speed** = distance / time (scalar)",),
             ("**Velocity** = displacement / time (vector)",),
             ("Describes motion over an **entire interval**",),
             ],
            size=14, gap=8)

c2 = add_rounded_rect(s, Inches(6.85), card_y, card_w, card_h, OFFWHITE, line=True, line_color=LIGHT_GRAY, radius=0.08, shadow=True)
tb, tf = add_textbox(s, Inches(7.15), card_y + Inches(0.2), card_w - Inches(0.6), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "INSTANTANEOUS", 14, RGBColor(0xC7, 0x6A, 0x1B), bold=True)
bullet_list(s, Inches(7.15), card_y + Inches(0.62), card_w - Inches(0.6),
            [("Speed/velocity at **one exact moment** in time",),
             ("What a speedometer shows **right now**",),
             ("The limit of average velocity as Δt → 0",),
             ],
            size=14, gap=8)

formula_card(s, Inches(0.6), Inches(4.35), Inches(5.85), Inches(1.55),
             "Average Velocity",
             "v_avg = Δx / Δt",
             note="Δx = displacement,  Δt = time elapsed", accent=ACCENT, bg=CARD_BLUE)

formula_card(s, Inches(6.85), Inches(4.35), Inches(5.85), Inches(1.55),
             "Average Speed",
             "speed_avg = d / t",
             note="d = total distance,  t = total time", accent=GOLD, bg=RGBColor(0xFC, 0xF2, 0xE3))

tip = add_rounded_rect(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.65), NAVY, radius=0.15)
tb, tf = add_textbox(s, Inches(0.95), Inches(6.15), Inches(11.4), Inches(0.65), anchor=MSO_ANCHOR.MIDDLE)
p = add_para(tf, first=True)
set_run(p.add_run(), "Units: ", 13.5, ACCENT, bold=True)
set_run(p.add_run(), "meters per second (m/s)   —   Velocity carries a sign; speed is always positive.", 13.5, WHITE)

footer(s, 5)


# ---------------------------------------------------------------------------
# SLIDE 6 — Acceleration
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Core Concepts", "Acceleration", 6)

section_label(s, Inches(0.6), Inches(1.65), "Definition")
bullet_list(s, Inches(0.6), Inches(2.1), Inches(5.9),
            [("The **rate of change of velocity** over time",),
             ("A **vector** — has direction, just like velocity",),
             ("Occurs whenever speed OR direction changes",),
             ],
            size=15.5, gap=12)

formula_card(s, Inches(0.6), Inches(4.2), Inches(5.9), Inches(1.85),
             "Average Acceleration",
             "a_avg = Δv / Δt",
             note="Δv = v_f − v_i     |     Units: meters per second² (m/s²)",
             accent=ACCENT, bg=CARD_BLUE)

# Right side: speeding up vs slowing down
c_y = Inches(1.65)
c_h = Inches(2.15)
c1 = add_rounded_rect(s, Inches(6.85), c_y, Inches(5.85), c_h, RGBColor(0xEC, 0xF7, 0xEE), line=True, line_color=RGBColor(0x3F, 0xA1, 0x5E), radius=0.08)
tb, tf = add_textbox(s, Inches(7.15), c_y + Inches(0.18), Inches(5.3), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "↑ SPEEDING UP", 15, RGBColor(0x2E, 0x7D, 0x46), bold=True)
bullet_list(s, Inches(7.15), c_y + Inches(0.65), Inches(5.3),
            [("Velocity and acceleration point the **same direction**",),
             ("e.g. v = +5 m/s, a = +2 m/s² → speeding up (+)",),
             ], size=14, gap=8)

c2_y = Inches(3.95)
c2 = add_rounded_rect(s, Inches(6.85), c2_y, Inches(5.85), c_h, RGBColor(0xFC, 0xEC, 0xEA), line=True, line_color=RED_ACCENT, radius=0.08)
tb, tf = add_textbox(s, Inches(7.15), c2_y + Inches(0.18), Inches(5.3), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "↓ SLOWING DOWN", 15, RED_ACCENT, bold=True)
bullet_list(s, Inches(7.15), c2_y + Inches(0.65), Inches(5.3),
            [("Velocity and acceleration point in **opposite directions**",),
             ("e.g. v = +5 m/s, a = −2 m/s² → slowing down (decelerating)",),
             ], size=14, gap=8)

note = add_rounded_rect(s, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.6), LIGHT_GRAY, radius=0.15)
tb, tf = add_textbox(s, Inches(0.95), Inches(6.25), Inches(11.4), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
p = add_para(tf, first=True)
set_run(p.add_run(), "Key idea:  ", 13, NAVY, bold=True)
set_run(p.add_run(), "The sign of acceleration alone doesn't tell you speeding up or slowing down — compare it to the sign of velocity.", 13, TEXT_BODY)

footer(s, 6)


# ---------------------------------------------------------------------------
# SLIDE 7 — The Big 4 Kinematic Equations
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Formula Toolkit", "The Big 4 Kinematic Equations", 7)

tb, tf = add_textbox(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.4))
p = add_para(tf, first=True)
set_run(p.add_run(), "Valid whenever acceleration is ", 14.5, TEXT_BODY, italic=True)
set_run(p.add_run(), "constant.", 14.5, TEXT_BODY, bold=True, italic=True)
set_run(p.add_run(), "  Variables: x = position, v = velocity, a = acceleration, t = time", 14.5, TEXT_BODY, italic=True)

eqs = [
    ("1", "v_f = v_i + a t", "Velocity as a function of time", "no Δx"),
    ("2", "Δx = v_i t + ½ a t²", "Position as a function of time", "no v_f"),
    ("3", "v_f² = v_i² + 2aΔx", "Velocity independent of time", "no t"),
    ("4", "Δx = ½ (v_i + v_f) t", "Average of initial & final velocity", "no a"),
]

grid_x = Inches(0.6)
grid_y = Inches(2.15)
gap = Inches(0.25)
cell_w = Inches((12.1 - 0.25) / 2)
cell_h = Inches(2.1)

positions = [
    (grid_x, grid_y),
    (grid_x + cell_w + gap, grid_y),
    (grid_x, grid_y + cell_h + gap),
    (grid_x + cell_w + gap, grid_y + cell_h + gap),
]

for (num, formula, desc, omits), (px, py) in zip(eqs, positions):
    card = add_rounded_rect(s, px, py, cell_w, cell_h, CARD_BLUE, line=True, line_color=ACCENT, line_w=Pt(1.25), radius=0.09, shadow=True)
    # number badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, px + Inches(0.25), py + Inches(0.22), Inches(0.42), Inches(0.42))
    badge.fill.solid(); badge.fill.fore_color.rgb = NAVY; badge.line.fill.background(); badge.shadow.inherit = False
    btf = badge.text_frame
    btf.margin_left = 0; btf.margin_top = 0
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    set_run(bp.add_run(), num, 16, WHITE, bold=True)

    tb, tf = add_textbox(s, px + Inches(0.85), py + Inches(0.22), cell_w - Inches(1.9), Inches(0.5))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = add_para(tf, first=True)
    set_run(p.add_run(), formula, 22, NAVY, bold=True, font="Cambria")

    pill(s, px + cell_w - Inches(1.15), py + Inches(0.24), Inches(0.9), Inches(0.36), omits, LIGHT_GRAY, TEXT_BODY, size=11)

    tb2, tf2 = add_textbox(s, px + Inches(0.25), py + Inches(0.95), cell_w - Inches(0.5), Inches(1.0))
    p2 = add_para(tf2, first=True, line_spacing=1.1)
    set_run(p2.add_run(), desc, 14, TEXT_BODY)

footer(s, 7)


# ---------------------------------------------------------------------------
# SLIDE 8 — Graphing Motion
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Visualizing Motion", "Graphing Motion", 8)

def mini_graph(slide, x, y, w, h, title, y_label, kind):
    card = add_rounded_rect(slide, x, y, w, h, OFFWHITE, line=True, line_color=LIGHT_GRAY, radius=0.06, shadow=True)
    pad_l = Inches(0.55)
    pad_b = Inches(0.4)
    tb, tf = add_textbox(slide, x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), Inches(0.32))
    p = add_para(tf, first=True)
    set_run(p.add_run(), title, 13.5, NAVY, bold=True)

    axis_x0 = x + pad_l
    axis_y0 = y + h - pad_b
    axis_w = w - pad_l - Inches(0.25)
    axis_h = h - pad_b - Inches(0.55)

    # axes
    xaxis = slide.shapes.add_connector(1, axis_x0, axis_y0, axis_x0 + axis_w, axis_y0)
    xaxis.line.color.rgb = MID_GRAY; xaxis.line.width = Pt(1.25)
    yaxis = slide.shapes.add_connector(1, axis_x0, axis_y0, axis_x0, axis_y0 - axis_h)
    yaxis.line.color.rgb = MID_GRAY; yaxis.line.width = Pt(1.25)

    # axis labels
    tbx, tfx = add_textbox(slide, axis_x0 + axis_w - Inches(0.5), axis_y0 + Inches(0.03), Inches(0.6), Inches(0.25))
    px = add_para(tfx, first=True)
    set_run(px.add_run(), "t", 11, TEXT_BODY, italic=True)
    tby, tfy = add_textbox(slide, x + Inches(0.05), y + Inches(0.4), Inches(0.5), Inches(0.3))
    py_ = add_para(tfy, first=True)
    set_run(py_.add_run(), y_label, 11, TEXT_BODY, italic=True)

    if kind == "line_up":
        conn = slide.shapes.add_connector(1, axis_x0, axis_y0, axis_x0 + axis_w, axis_y0 - axis_h)
        conn.line.color.rgb = ACCENT_DARK; conn.line.width = Pt(2.5)
        cap = "Constant + slope → steady increase"
    elif kind == "curve_up":
        # approximate curve with freeform-like multiple segments (parabola opening upward via connectors)
        n = 12
        pts = []
        for i in range(n + 1):
            t = i / n
            xx = axis_x0 + Emu(int(axis_w * t))
            yy = axis_y0 - Emu(int(axis_h * (t ** 2)))
            pts.append((xx, yy))
        for i in range(len(pts) - 1):
            seg = slide.shapes.add_connector(1, pts[i][0], pts[i][1], pts[i+1][0], pts[i+1][1])
            seg.line.color.rgb = GOLD; seg.line.width = Pt(2.5)
        cap = "Curving slope → changing velocity"
    elif kind == "flat_step":
        mid_x = axis_x0 + Emu(int(axis_w * 0.55))
        mid_y = axis_y0 - Emu(int(axis_h * 0.6))
        seg1 = slide.shapes.add_connector(1, axis_x0, axis_y0, mid_x, mid_y)
        seg1.line.color.rgb = ACCENT_DARK; seg1.line.width = Pt(2.5)
        seg2 = slide.shapes.add_connector(1, mid_x, mid_y, axis_x0 + axis_w, mid_y)
        seg2.line.color.rgb = ACCENT_DARK; seg2.line.width = Pt(2.5)
        cap = "Flat segment → zero slope (at rest)"
    elif kind == "const_pos":
        yy = axis_y0 - Emu(int(axis_h * 0.55))
        seg = slide.shapes.add_connector(1, axis_x0, yy, axis_x0 + axis_w, yy)
        seg.line.color.rgb = GOLD; seg.line.width = Pt(2.5)
        cap = "Horizontal line → constant velocity"
    elif kind == "const_flat_zero":
        seg = slide.shapes.add_connector(1, axis_x0, axis_y0, axis_x0 + axis_w, axis_y0)
        seg.line.color.rgb = RED_ACCENT; seg.line.width = Pt(2.5)
        cap = "On axis → zero acceleration"

    tbc, tfc = add_textbox(slide, x + Inches(0.2), y + h - Inches(0.28), w - Inches(0.4), Inches(0.3))
    pc = add_para(tfc, first=True)
    set_run(pc.add_run(), cap, 10.5, MID_GRAY, italic=True)


g_y = Inches(1.65)
g_h = Inches(2.55)
g_w = Inches(3.85)
gap = Inches(0.28)
mini_graph(s, Inches(0.6), g_y, g_w, g_h, "Position vs. Time (x–t)", "x", "curve_up")
mini_graph(s, Inches(0.6) + g_w + gap, g_y, g_w, g_h, "Velocity vs. Time (v–t)", "v", "const_pos")
mini_graph(s, Inches(0.6) + 2*(g_w + gap), g_y, g_w, g_h, "Acceleration vs. Time (a–t)", "a", "const_flat_zero")

# slope/area explainer
tip_y = Inches(4.5)
tip1 = add_rounded_rect(s, Inches(0.6), tip_y, Inches(5.85), Inches(2.35), CARD_BLUE, line=True, line_color=ACCENT, radius=0.08)
tb, tf = add_textbox(s, Inches(0.9), tip_y + Inches(0.18), Inches(5.3), Inches(0.35))
p = add_para(tf, first=True)
set_run(p.add_run(), "SLOPE reveals the next quantity", 14.5, ACCENT_DARK, bold=True)
bullet_list(s, Inches(0.9), tip_y + Inches(0.65), Inches(5.3),
            [("Slope of **x–t graph** = velocity",),
             ("Slope of **v–t graph** = acceleration",),
             ], size=14, gap=10)

tip2 = add_rounded_rect(s, Inches(6.85), tip_y, Inches(5.85), Inches(2.35), RGBColor(0xFC, 0xF2, 0xE3), line=True, line_color=GOLD, radius=0.08)
tb, tf = add_textbox(s, Inches(7.15), tip_y + Inches(0.18), Inches(5.3), Inches(0.35))
p = add_para(tf, first=True)
set_run(p.add_run(), "AREA reveals the previous quantity", 14.5, RGBColor(0xC7, 0x6A, 0x1B), bold=True)
bullet_list(s, Inches(7.15), tip_y + Inches(0.65), Inches(5.3),
            [("Area under **v–t graph** = displacement",),
             ("Area under **a–t graph** = change in velocity",),
             ], size=14, gap=10)

footer(s, 8)


# ---------------------------------------------------------------------------
# SLIDE 9 — Free Fall & Gravity
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Special Case", "Free Fall & Gravity", 9)

section_label(s, Inches(0.6), Inches(1.65), "What is Free Fall?")
bullet_list(s, Inches(0.6), Inches(2.1), Inches(5.9),
            [("Motion under the influence of **gravity alone** — no air resistance",),
             ("All objects fall with the **same acceleration**, regardless of mass",),
             ("Acceleration is always directed **downward**, toward Earth's center",),
             ],
            size=15.5, gap=12)

formula_card(s, Inches(0.6), Inches(4.55), Inches(5.9), Inches(1.65),
             "Acceleration Due to Gravity",
             "g ≈ 9.8 m/s²  (downward)",
             note="Some texts round to 10 m/s² for quick estimates", accent=RED_ACCENT,
             bg=RGBColor(0xFC, 0xEC, 0xEA))

# Right visual: falling object with arrow
vis = add_rounded_rect(s, Inches(6.85), Inches(1.65), Inches(5.85), Inches(4.55), NAVY, radius=0.06)
tb, tf = add_textbox(s, Inches(7.15), Inches(1.85), Inches(5.3), Inches(0.35))
p = add_para(tf, first=True)
set_run(p.add_run(), "Sign convention (common choice)", 13.5, ACCENT, bold=True)

ball = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.55), Inches(2.35), Inches(0.5), Inches(0.5))
ball.fill.solid(); ball.fill.fore_color.rgb = GOLD; ball.line.fill.background(); ball.shadow.inherit = False

arrow2 = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.68), Inches(3.0), Inches(0.25), Inches(2.4))
arrow2.fill.solid(); arrow2.fill.fore_color.rgb = ACCENT; arrow2.line.fill.background(); arrow2.shadow.inherit = False

tb, tf = add_textbox(s, Inches(10.05), Inches(3.6), Inches(2.3), Inches(0.8))
p = add_para(tf, first=True, line_spacing=1.15)
set_run(p.add_run(), "g = −9.8 m/s²", 16, WHITE, bold=True)
p2 = add_para(tf, line_spacing=1.15)
set_run(p2.add_run(), "(if up is +)", 12.5, LIGHT_GRAY, italic=True)

bullet_pts_y = Inches(5.6)
tb, tf = add_textbox(s, Inches(7.15), bullet_pts_y, Inches(5.3), Inches(0.5))
p = add_para(tf, first=True, line_spacing=1.2)
set_run(p.add_run(), "• Up = positive, Down = negative → g is negative", 13, LIGHT_GRAY)
p2 = add_para(tf, line_spacing=1.2)
set_run(p2.add_run(), "• The Big 4 equations apply directly, using a = g", 13, LIGHT_GRAY)

footer(s, 9)


# ---------------------------------------------------------------------------
# SLIDE 10 — Problem-Solving Steps & Summary
# ---------------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "Wrap-Up", "Problem-Solving Steps & Summary", 10)

section_label(s, Inches(0.6), Inches(1.6), "4-Step Problem-Solving Method")

steps = [
    ("1", "GIVEN", "List every known value from the problem (with units & signs)"),
    ("2", "UNKNOWN", "Identify exactly what quantity you're solving for"),
    ("3", "EQUATION", "Pick the Big 4 equation that has your knowns + unknown (no extras)"),
    ("4", "SOLVE", "Substitute values, solve algebraically, and check units"),
]
sx = Inches(0.6)
sw = Inches(2.85)
sgap = Inches(0.2)
sy = Inches(2.1)
sh = Inches(2.0)
colors = [ACCENT, GOLD, RGBColor(0x7A, 0x6F, 0xDB), RED_ACCENT]
for i, ((num, label, desc), c) in enumerate(zip(steps, colors)):
    px = sx + i * (sw + sgap)
    card = add_rounded_rect(s, px, sy, sw, sh, OFFWHITE, line=True, line_color=c, line_w=Pt(1.5), radius=0.1, shadow=True)
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL, px + Inches(0.2), sy + Inches(0.2), Inches(0.45), Inches(0.45))
    badge.fill.solid(); badge.fill.fore_color.rgb = c; badge.line.fill.background(); badge.shadow.inherit = False
    btf = badge.text_frame; btf.margin_left = 0
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    set_run(bp.add_run(), num, 16, WHITE, bold=True)
    tb, tf = add_textbox(s, px + Inches(0.2), sy + Inches(0.8), sw - Inches(0.4), Inches(0.3))
    p = add_para(tf, first=True)
    set_run(p.add_run(), label, 14, NAVY, bold=True)
    tb2, tf2 = add_textbox(s, px + Inches(0.2), sy + Inches(1.15), sw - Inches(0.4), Inches(0.8))
    p2 = add_para(tf2, first=True, line_spacing=1.05)
    set_run(p2.add_run(), desc, 11.5, TEXT_BODY)
    if i < 3:
        arr = s.shapes.add_shape(MSO_SHAPE.CHEVRON, px + sw + Inches(0.01), sy + sh/2 - Inches(0.12), Inches(0.18), Inches(0.24))
        arr.fill.solid(); arr.fill.fore_color.rgb = MID_GRAY; arr.line.fill.background(); arr.shadow.inherit = False

# Summary recap
section_label(s, Inches(0.6), Inches(4.4), "Unit Recap")
recap = add_rounded_rect(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(1.85), NAVY, radius=0.08)
recap_items = [
    ("x, Δx", "Position & Displacement"),
    ("v_avg = Δx/Δt", "Velocity"),
    ("a_avg = Δv/Δt", "Acceleration"),
    ("Big 4 Equations", "Constant acceleration"),
    ("g ≈ 9.8 m/s²", "Free fall"),
]
rx = Inches(0.9)
rw = Inches(2.28)
for i, (formula, label) in enumerate(recap_items):
    px = rx + i * rw
    tb, tf = add_textbox(s, px, Inches(5.1), rw - Inches(0.15), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    p = add_para(tf, first=True, align=PP_ALIGN.LEFT, line_spacing=1.1)
    set_run(p.add_run(), formula, 15.5, ACCENT, bold=True, font="Cambria")
    tb2, tf2 = add_textbox(s, px, Inches(5.75), rw - Inches(0.15), Inches(0.7))
    p2 = add_para(tf2, first=True, line_spacing=1.1)
    set_run(p2.add_run(), label, 12, LIGHT_GRAY)
    if i < len(recap_items) - 1:
        add_rect(s, px + rw - Inches(0.15), Inches(5.05), Pt(1), Inches(1.4), NAVY_LIGHT)

tb, tf = add_textbox(s, Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4))
p = add_para(tf, first=True, align=PP_ALIGN.CENTER)
set_run(p.add_run(), "Next up: Unit 2 — Kinematics in Two Dimensions (Projectile Motion)", 13, MID_GRAY, italic=True)


# ---------------------------------------------------------------------------
prs.save("1D_Kinematics_Intro.pptx")
print("Saved 1D_Kinematics_Intro.pptx")
